from fastapi import APIRouter, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel
from typing import Optional, List, Literal
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

router = APIRouter()

BRAND_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BRAND_BLUE = RGBColor(0x4C, 0x6E, 0xF5)
BRAND_VIOLET = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xF0)
MUTED = RGBColor(0x88, 0x88, 0x99)


class Slide(BaseModel):
    title: str
    content: Optional[str] = ""
    bullets: Optional[List[str]] = None
    slideType: Literal["title", "content", "bullets", "conclusion"] = "content"


class PPTXRequest(BaseModel):
    title: str
    subtitle: Optional[str] = None
    slides: List[Slide]


def set_bg_color(slide, r, g, b):
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_text_box(slide, text, left, top, width, height, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


@router.post("/pptx")
async def generate_pptx(req: PPTXRequest):
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        W = prs.slide_width
        H = prs.slide_height
        blank_layout = prs.slide_layouts[6]  # completely blank

        for i, slide_data in enumerate(req.slides):
            slide = prs.slides.add_slide(blank_layout)

            # Background: dark
            set_bg_color(slide, 0x0F, 0x0F, 0x1A)

            # Accent bar on left
            bar = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(0), Inches(0),
                Inches(0.12), H
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = BRAND_BLUE
            bar.line.fill.background()

            is_title_slide = slide_data.slideType == "title" or i == 0
            is_conclusion = slide_data.slideType == "conclusion"

            if is_title_slide:
                # Large centered title
                add_text_box(
                    slide, req.title,
                    Inches(1), Inches(2),
                    Inches(11.3), Inches(1.5),
                    font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER
                )
                if req.subtitle:
                    add_text_box(
                        slide, req.subtitle,
                        Inches(1), Inches(3.8),
                        Inches(11.3), Inches(0.6),
                        font_size=20, color=MUTED, align=PP_ALIGN.CENTER
                    )
                # Slide number: 1
                add_text_box(slide, f"DocGenius AI", Inches(1), Inches(6.8), Inches(5), Inches(0.4), font_size=10, color=MUTED)
            else:
                # Slide number badge
                n_box = slide.shapes.add_shape(1, Inches(11.8), Inches(6.8), Inches(1.3), Inches(0.45))
                n_box.fill.solid()
                n_box.fill.fore_color.rgb = BRAND_BLUE
                n_box.line.fill.background()
                n_tf = n_box.text_frame.paragraphs[0]
                n_tf.alignment = PP_ALIGN.CENTER
                n_run = n_tf.add_run()
                n_run.text = f"{i + 1} / {len(req.slides)}"
                n_run.font.size = Pt(9)
                n_run.font.color.rgb = WHITE

                # Section heading chip
                if is_conclusion:
                    chip = slide.shapes.add_shape(1, Inches(0.2), Inches(0.2), Inches(2), Inches(0.38))
                    chip.fill.solid()
                    chip.fill.fore_color.rgb = BRAND_VIOLET
                    chip.line.fill.background()
                    chip_p = chip.text_frame.paragraphs[0]
                    chip_p.alignment = PP_ALIGN.CENTER
                    chip_r = chip_p.add_run()
                    chip_r.text = "CONCLUSION"
                    chip_r.font.size = Pt(9)
                    chip_r.font.color.rgb = WHITE

                # Title
                add_text_box(
                    slide, slide_data.title,
                    Inches(0.3), Inches(0.3),
                    Inches(12.5), Inches(1.0),
                    font_size=26, bold=True, color=WHITE
                )

                # Divider
                line = slide.shapes.add_shape(1, Inches(0.3), Inches(1.4), Inches(6), Inches(0.04))
                line.fill.solid()
                line.fill.fore_color.rgb = BRAND_BLUE
                line.line.fill.background()

                # Content body
                if slide_data.content:
                    add_text_box(
                        slide, slide_data.content,
                        Inches(0.3), Inches(1.6),
                        Inches(12.7), Inches(2.0),
                        font_size=15, color=LIGHT_GRAY
                    )

                # Bullets
                if slide_data.bullets:
                    top_offset = Inches(3.8) if slide_data.content else Inches(1.8)
                    for j, bullet in enumerate(slide_data.bullets[:5]):
                        bullet_y = top_offset + j * Inches(0.55)
                        dot = slide.shapes.add_shape(1, Inches(0.3), bullet_y + Inches(0.12), Inches(0.08), Inches(0.08))
                        dot.fill.solid()
                        dot.fill.fore_color.rgb = BRAND_BLUE
                        dot.line.fill.background()
                        add_text_box(
                            slide, bullet,
                            Inches(0.55), bullet_y,
                            Inches(12.3), Inches(0.5),
                            font_size=14, color=LIGHT_GRAY
                        )

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        return Response(
            content=buffer.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{req.title}.pptx"'},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
