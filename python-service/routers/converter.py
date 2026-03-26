from fastapi import APIRouter, HTTPException, UploadFile, File, Form
from fastapi.responses import Response
from typing import List, Tuple
from io import BytesIO
import zipfile
import os

import fitz
from pypdf import PdfReader
from docx import Document as WordDocument
from pptx import Presentation
from openpyxl import Workbook, load_workbook
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

router = APIRouter()


def split_paragraphs(text: str) -> List[str]:
    lines = [line.strip() for line in text.replace("\r", "").split("\n")]
    chunks: List[str] = []
    current: List[str] = []

    for line in lines:
        if not line:
            if current:
                chunks.append(" ".join(current).strip())
                current = []
            continue
        current.append(line)

    if current:
        chunks.append(" ".join(current).strip())

    return chunks[:24] if chunks else ["No readable content found."]


def extract_text_from_pdf(data: bytes) -> str:
    try:
        document = fitz.open(stream=data, filetype="pdf")
        pages: List[str] = []
        for page in document:
            text = page.get_text("text").strip()
            if text:
                pages.append(text)

        extracted = "\n\n".join(pages).strip()
        if extracted:
            return extracted
    except Exception:
        pass

    reader = PdfReader(BytesIO(data))
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    text = "\n\n".join(page for page in pages if page)
    return text.strip() or "No readable PDF text found."


def extract_text_from_docx(data: bytes) -> str:
    doc = WordDocument(BytesIO(data))
    parts = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
    return "\n\n".join(parts) or "No readable DOCX text found."


def extract_text_from_pptx(data: bytes) -> str:
    prs = Presentation(BytesIO(data))
    parts: List[str] = []
    for slide in prs.slides:
      for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
          parts.append(shape.text.strip())
    return "\n\n".join(parts) or "No readable PPTX text found."


def extract_text_from_xlsx(data: bytes) -> str:
    wb = load_workbook(BytesIO(data), data_only=True)
    lines: List[str] = []
    for ws in wb.worksheets:
        lines.append(f"Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines) or "No readable XLSX text found."


def infer_source_type(filename: str) -> str:
    return os.path.splitext(filename)[1].lower().replace(".", "")


def build_docx(title: str, text: str) -> Tuple[str, bytes, str]:
    doc = WordDocument()
    doc.add_heading(title, 0)
    for chunk in split_paragraphs(text):
        doc.add_paragraph(chunk)
    buffer = BytesIO()
    doc.save(buffer)
    return f"{title}.docx", buffer.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def build_pdf(title: str, text: str) -> Tuple[str, bytes, str]:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, title=title)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    for chunk in split_paragraphs(text):
        story.extend([Paragraph(chunk, styles["BodyText"]), Spacer(1, 8)])
    doc.build(story)
    return f"{title}.pdf", buffer.getvalue(), "application/pdf"


def build_pptx(title: str, text: str) -> Tuple[str, bytes, str]:
    prs = Presentation()
    chunks = split_paragraphs(text)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = "Converted by DocGenius AI"

    for chunk in chunks[:10]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = chunk[:60] if len(chunk) > 60 else chunk
        slide.placeholders[1].text = chunk

    buffer = BytesIO()
    prs.save(buffer)
    return f"{title}.pptx", buffer.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def build_xlsx(title: str, text: str) -> Tuple[str, bytes, str]:
    wb = Workbook()
    ws = wb.active
    ws.title = "Converted Data"
    ws.append(["Line", "Content"])
    for index, chunk in enumerate(split_paragraphs(text), start=1):
        ws.append([index, chunk])
    buffer = BytesIO()
    wb.save(buffer)
    return f"{title}.xlsx", buffer.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def build_txt(title: str, text: str) -> Tuple[str, bytes, str]:
    return f"{title}.txt", text.encode("utf-8"), "text/plain; charset=utf-8"


def build_pdf_images(title: str, data: bytes) -> Tuple[str, bytes, str]:
    document = fitz.open(stream=data, filetype="pdf")
    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for index, page in enumerate(document, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            archive.writestr(f"{title}-page-{index}.png", pix.tobytes("png"))
    return f"{title}-images.zip", zip_buffer.getvalue(), "application/zip"


@router.post("/convert")
async def convert_document(
    file: UploadFile = File(...),
    target_type: str = Form(...),
):
    try:
        source_type = infer_source_type(file.filename or "document")
        target_type = target_type.lower().strip()
        data = await file.read()
        title = os.path.splitext(file.filename or "document")[0][:40] or "document"

        if source_type == "pdf" and target_type == "img":
            file_name, payload, media_type = build_pdf_images(title, data)
            return Response(content=payload, media_type=media_type, headers={"Content-Disposition": f'attachment; filename="{file_name}"'})

        if source_type == target_type:
            return Response(content=data, media_type=file.content_type or "application/octet-stream")

        if source_type == "pdf":
            extracted = extract_text_from_pdf(data)
        elif source_type == "docx":
            extracted = extract_text_from_docx(data)
        elif source_type == "pptx":
            extracted = extract_text_from_pptx(data)
        elif source_type == "xlsx":
            extracted = extract_text_from_xlsx(data)
        elif source_type == "txt":
            extracted = data.decode("utf-8", errors="ignore").strip() or "No readable text found."
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported source type: {source_type}")

        builders = {
            "docx": build_docx,
            "pdf": build_pdf,
            "pptx": build_pptx,
            "xlsx": build_xlsx,
            "txt": build_txt,
        }

        if target_type not in builders:
            raise HTTPException(status_code=400, detail=f"Unsupported target type: {target_type}")

        file_name, payload, media_type = builders[target_type](title, extracted)
        return Response(
            content=payload,
            media_type=media_type,
            headers={"Content-Disposition": f'attachment; filename="{file_name}"'},
        )
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))
